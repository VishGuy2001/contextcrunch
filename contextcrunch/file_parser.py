"""
file_parser.py — PDF, PPTX, DOCX, XLSX, images, code
"""
import os
import io
from dataclasses import dataclass, field
from typing import Optional
from .tokenizer import count_tokens, count_code_tokens, image_tokens, detect_language


@dataclass
class ParseResult:
    text: str
    token_estimate: int
    file_type: str
    filename: str
    breakdown: dict = field(default_factory=dict)
    images_found: int = 0
    image_token_estimate: int = 0
    language: Optional[str] = None
    pages: Optional[int] = None
    slides: Optional[list] = None
    warning: Optional[str] = None


def parse_file(file_bytes: bytes, filename: str, model: str = "claude", plan: str = "plus") -> ParseResult:
    ext = os.path.splitext(filename)[1].lower()
    code_exts = {".py",".js",".ts",".java",".cpp",".c",".rs",".go",".sql",".html",".css",".json",".yaml",".yml"}
    parsers = {".txt": _text, ".md": _text, ".csv": _text, ".pdf": _pdf, ".docx": _docx, ".doc": _docx, ".pptx": _pptx, ".ppt": _pptx, ".xlsx": _xlsx, ".xls": _xlsx, ".png": _image, ".jpg": _image, ".jpeg": _image, ".webp": _image, ".gif": _image}
    if ext in code_exts:
        return _code(file_bytes, filename, model, plan)
    parser = parsers.get(ext, _text)
    return parser(file_bytes, filename, model, plan)


def _text(fb, fn, model, plan):
    try:
        text = fb.decode("utf-8", errors="replace")
        return ParseResult(text=text, token_estimate=count_tokens(text, model, plan), file_type="text", filename=fn, breakdown={"tokens": count_tokens(text, model, plan)})
    except Exception as e:
        return ParseResult(text="", token_estimate=0, file_type="text", filename=fn, warning=str(e))


def _pdf(fb, fn, model, plan):
    try:
        import fitz
        doc = fitz.open(stream=fb, filetype="pdf")
        all_text, imgs, img_tokens = [], 0, 0
        for page in doc:
            all_text.append(page.get_text())
            for img in page.get_images():
                imgs += 1
                try:
                    pix = fitz.Pixmap(doc, img[0])
                    img_tokens += image_tokens(pix.width, pix.height, model)
                except:
                    img_tokens += image_tokens(512, 512, model)
        full = "\n".join(all_text)
        txt_tok = count_tokens(full, model, plan)
        return ParseResult(text=full, token_estimate=txt_tok+img_tokens, file_type="pdf", filename=fn, pages=len(doc), images_found=imgs, image_token_estimate=img_tokens, breakdown={"text": txt_tok, "images": img_tokens, "total": txt_tok+img_tokens}, warning=f"{imgs} image(s) — token count estimated." if imgs else None)
    except ImportError:
        return ParseResult(text="", token_estimate=0, file_type="pdf", filename=fn, warning="PyMuPDF not installed: pip install pymupdf")
    except Exception as e:
        return ParseResult(text="", token_estimate=0, file_type="pdf", filename=fn, warning=str(e))


def _docx(fb, fn, model, plan):
    try:
        from docx import Document
        doc = Document(io.BytesIO(fb))
        paras = [p.text for p in doc.paragraphs if p.text.strip()]
        text = "\n".join(paras)
        return ParseResult(text=text, token_estimate=count_tokens(text, model, plan), file_type="docx", filename=fn, breakdown={"paragraphs": len(paras), "tokens": count_tokens(text, model, plan)})
    except ImportError:
        return ParseResult(text="", token_estimate=0, file_type="docx", filename=fn, warning="python-docx not installed")
    except Exception as e:
        return ParseResult(text="", token_estimate=0, file_type="docx", filename=fn, warning=str(e))


def _pptx(fb, fn, model, plan):
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(fb))
        slides_data, all_text, total_img_tok, total_imgs = [], [], 0, 0
        for i, slide in enumerate(prs.slides):
            slide_text, slide_imgs, slide_img_tok = [], 0, 0
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_text.append(para.text.strip())
                if shape.shape_type == 13:
                    total_imgs += 1; slide_imgs += 1
                    try:
                        w, h = int(shape.width/9144), int(shape.height/9144)
                        t = image_tokens(max(w,100), max(h,100), model)
                    except:
                        t = image_tokens(512, 512, model)
                    slide_img_tok += t; total_img_tok += t
            st = " ".join(slide_text)
            stok = count_tokens(st, model, plan)
            all_text.append(st)
            slides_data.append({"slide": i+1, "text_tokens": stok, "image_tokens": slide_img_tok, "images": slide_imgs, "total_tokens": stok+slide_img_tok, "preview": st[:80]+"..." if len(st)>80 else st})
        full = "\n".join(all_text)
        txt_tok = count_tokens(full, model, plan)
        return ParseResult(text=full, token_estimate=txt_tok+total_img_tok, file_type="pptx", filename=fn, slides=slides_data, images_found=total_imgs, image_token_estimate=total_img_tok, breakdown={"slides": len(prs.slides), "text_tokens": txt_tok, "image_tokens": total_img_tok, "total": txt_tok+total_img_tok})
    except ImportError:
        return ParseResult(text="", token_estimate=0, file_type="pptx", filename=fn, warning="python-pptx not installed")
    except Exception as e:
        return ParseResult(text="", token_estimate=0, file_type="pptx", filename=fn, warning=str(e))


def _xlsx(fb, fn, model, plan):
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(fb), read_only=True)
        rows = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                rt = " ".join(str(c) for c in row if c is not None)
                if rt.strip():
                    rows.append(rt)
        text = "\n".join(rows)
        return ParseResult(text=text, token_estimate=count_tokens(text, model, plan), file_type="xlsx", filename=fn, breakdown={"sheets": len(wb.worksheets), "rows": len(rows), "tokens": count_tokens(text, model, plan)})
    except ImportError:
        return ParseResult(text="", token_estimate=0, file_type="xlsx", filename=fn, warning="openpyxl not installed")
    except Exception as e:
        return ParseResult(text="", token_estimate=0, file_type="xlsx", filename=fn, warning=str(e))


def _image(fb, fn, model, plan):
    try:
        from PIL import Image
        img = Image.open(io.BytesIO(fb))
        w, h = img.size
        tok = image_tokens(w, h, model)
        return ParseResult(text=f"[Image: {fn}, {w}x{h}px]", token_estimate=tok, file_type="image", filename=fn, images_found=1, image_token_estimate=tok, breakdown={"width": w, "height": h, "tokens": tok, "formula": {"claude": f"({w}×{h})/750={int((w*h)/750)}", "chatgpt": f"85+tiles×170={85+((w+511)//512)*((h+511)//512)*170}", "gemini": "258 fixed"}})
    except Exception:
        tok = image_tokens(512, 512, model)
        return ParseResult(text=f"[Image: {fn}]", token_estimate=tok, file_type="image", filename=fn, images_found=1, image_token_estimate=tok, warning="Could not read image dimensions — using 512×512 estimate")


def _code(fb, fn, model, plan):
    try:
        text = fb.decode("utf-8", errors="replace")
        lang = detect_language(fn)
        tok = count_code_tokens(text, lang)
        lines = text.split("\n")
        return ParseResult(text=text, token_estimate=tok, file_type="code", filename=fn, language=lang, breakdown={"language": lang, "lines": len(lines), "tokens": tok})
    except Exception as e:
        return ParseResult(text="", token_estimate=0, file_type="code", filename=fn, warning=str(e))
